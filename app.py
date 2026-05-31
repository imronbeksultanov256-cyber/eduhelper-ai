import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import os
import json
import re

# ─────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────
st.set_page_config(
    page_title="EduHelper AI",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─────────────────────────────────────────
# CSS
# ─────────────────────────────────────────
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(135deg, #1e3a5f 0%, #2d6a4f 100%);
        padding: 1.5rem 2rem;
        border-radius: 12px;
        margin-bottom: 1.5rem;
        color: white;
        text-align: center;
    }
    .main-header h1 { font-size: 2.2rem; margin: 0; }
    .main-header p  { font-size: 1rem; margin: 0.3rem 0 0; opacity: 0.85; }
    .module-card {
        background: #f8f9fa;
        border-left: 4px solid #2d6a4f;
        padding: 1rem 1.2rem;
        border-radius: 8px;
        margin-bottom: 1rem;
    }
    .result-box {
        background: #eef6f1;
        border: 1px solid #b7d5c4;
        border-radius: 8px;
        padding: 1.2rem;
        margin-top: 1rem;
        white-space: pre-wrap;
        font-size: 0.93rem;
        line-height: 1.6;
    }
    .stButton > button {
        background: #2d6a4f;
        color: white;
        border-radius: 8px;
        border: none;
        padding: 0.5rem 1.4rem;
        font-weight: 600;
    }
    .stButton > button:hover { background: #1e4d38; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────
# SIDEBAR
# ─────────────────────────────────────────
with st.sidebar:
    st.markdown("## 🎓 EduHelper AI")
    st.divider()

    role = st.radio("👤 Ваша роль", ["🎒 Студент", "👩‍🏫 Преподаватель"])

    st.divider()
    st.markdown("### 🔑 Gemini API Key")
    user_key = st.text_input(
        "Вставьте свой ключ (необязательно)",
        type="password",
        placeholder="AIza...",
        help="Бесплатный ключ: aistudio.google.com",
    )
    st.caption("💡 Получите бесплатный ключ на [Google AI Studio](https://aistudio.google.com/app/apikey)")
    st.divider()
    st.caption("Powered by Google Gemini 1.5 Flash · python-pptx · Streamlit")

# ─────────────────────────────────────────
# API INIT
# ─────────────────────────────────────────
def get_api_key() -> str | None:
    if user_key and user_key.strip():
        return user_key.strip()
    try:
        return st.secrets["GEMINI_API_KEY"]
    except Exception:
        pass
    return os.environ.get("GEMINI_API_KEY")

def get_model():
    key = get_api_key()
    if not key:
        return None, "❌ API-ключ не найден. Введите ключ в боковом меню или настройте Secrets."
    try:
        genai.configure(api_key=key)
        model = genai.GenerativeModel("gemini-1.5-flash")
        return model, None
    except Exception as e:
        return None, f"❌ Ошибка инициализации API: {e}"

def call_gemini(prompt: str) -> str:
    model, err = get_model()
    if err:
        return f"__ERROR__: {err}"
    try:
        resp = model.generate_content(prompt)
        return resp.text
    except Exception as e:
        return f"__ERROR__: Ошибка запроса к Gemini: {e}"

# ─────────────────────────────────────────
# HEADER
# ─────────────────────────────────────────
st.markdown("""
<div class="main-header">
  <h1>🎓 EduHelper AI</h1>
  <p>Бесплатный академический помощник на базе Google Gemini</p>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────
# PPTX BUILDER
# ─────────────────────────────────────────
def build_pptx(slides_data: list[dict]) -> bytes:
    """
    slides_data: [{"title": str, "bullets": [str, ...]}, ...]
    Returns bytes of .pptx file.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BG    = RGBColor(0x1E, 0x3A, 0x5F)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT     = RGBColor(0x2D, 0x6A, 0x4F)
    LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG if idx == 0 else LIGHT_GRAY

        # Accent bar (left)
        from pptx.util import Emu
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0),
            Inches(0.18), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # Title box
        title_tf = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.3),
            Inches(12.5), Inches(1.4),
        ).text_frame
        title_tf.word_wrap = True
        p = title_tf.paragraphs[0]
        p.text = slide_info.get("title", "")
        p.font.size = Pt(32 if idx == 0 else 26)
        p.font.bold = True
        p.font.color.rgb = WHITE if idx == 0 else DARK_BG

        # Bullets
        bullets = slide_info.get("bullets", [])
        if bullets:
            content_tf = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.9),
                Inches(12.3), Inches(5.2),
            ).text_frame
            content_tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    para = content_tf.paragraphs[0]
                else:
                    para = content_tf.add_paragraph()
                para.text = f"▸  {bullet}"
                para.font.size = Pt(18)
                para.font.color.rgb = WHITE if idx == 0 else RGBColor(0x2C, 0x2C, 0x2C)
                para.space_before = Pt(8)

        # Slide number (not on title slide)
        if idx > 0:
            num_tf = slide.shapes.add_textbox(
                Inches(12.6), Inches(7.1),
                Inches(0.7), Inches(0.35),
            ).text_frame
            p = num_tf.paragraphs[0]
            p.text = str(idx + 1)
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
            p.alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def parse_slides_json(raw: str) -> list[dict] | None:
    """Extract JSON array from model output."""
    # Try to find JSON array
    match = re.search(r'\[[\s\S]*\]', raw)
    if not match:
        return None
    try:
        data = json.loads(match.group())
        if isinstance(data, list):
            return data
    except Exception:
        pass
    return None

# ─────────────────────────────────────────
# ── СТУДЕНТ ──────────────────────────────
# ─────────────────────────────────────────
if "🎒 Студент" in role:
    tab1, tab2 = st.tabs(["📊 Ассистент Презентаций", "📝 ГОСТ-Корректор и Саммаризатор"])

    # ── TAB 1: ПРЕЗЕНТАЦИИ ─────────────────
    with tab1:
        st.markdown('<div class="module-card"><b>📊 Ассистент Презентаций</b><br>Введите тему или вставьте текст — ИИ построит структуру и создаст файл .pptx</div>', unsafe_allow_html=True)

        col1, col2 = st.columns([3, 1])
        with col1:
            pres_topic = st.text_area(
                "Тема или текст конспекта / реферата",
                height=180,
                placeholder="Например: «Клеточный цикл и митоз» или вставьте готовый текст...",
            )
        with col2:
            num_slides = st.slider("Кол-во слайдов", min_value=5, max_value=12, value=8)
            lang = st.selectbox("Язык", ["Русский", "English", "Uzbek"])
            generate_btn = st.button("🚀 Создать презентацию", use_container_width=True)

        if generate_btn:
            if not pres_topic.strip():
                st.warning("⚠️ Введите тему или текст для презентации.")
            else:
                with st.spinner("🤖 ИИ структурирует материал..."):
                    prompt = f"""You are a presentation specialist. Based on the input below, create a presentation structure with exactly {num_slides} slides.

Input (language: {lang}):
\"\"\"{pres_topic}\"\"\"

Return ONLY a valid JSON array (no markdown, no extra text) with this structure:
[
  {{"title": "Slide title", "bullets": ["Point 1", "Point 2", "Point 3"]}},
  ...
]

Rules:
- First slide: title slide with presentation name and 1-2 subtitle bullets
- Last slide: conclusion/summary slide
- Each slide: 3-4 short, informative bullet points (max 12 words each)
- All text must be in {lang}
- Total slides: exactly {num_slides}
"""
                    raw = call_gemini(prompt)

                if raw.startswith("__ERROR__"):
                    st.error(raw.replace("__ERROR__: ", ""))
                else:
                    slides_data = parse_slides_json(raw)
                    if not slides_data:
                        st.error("❌ Не удалось разобрать структуру от ИИ. Попробуйте ещё раз.")
                        with st.expander("Ответ модели (отладка)"):
                            st.code(raw)
                    else:
                        st.success(f"✅ Структура готова: {len(slides_data)} слайдов")
                        with st.expander("📋 Предпросмотр структуры"):
                            for i, s in enumerate(slides_data):
                                st.markdown(f"**{i+1}. {s.get('title','')}**")
                                for b in s.get("bullets", []):
                                    st.markdown(f"   • {b}")

                        with st.spinner("🎨 Сборка .pptx файла..."):
                            pptx_bytes = build_pptx(slides_data)

                        st.download_button(
                            label="⬇️ Скачать презентацию (.pptx)",
                            data=pptx_bytes,
                            file_name="presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )

    # ── TAB 2: ГОСТ-КОРРЕКТОР ──────────────
    with tab2:
        st.markdown('<div class="module-card"><b>📝 ГОСТ-Корректор и Саммаризатор</b><br>ИИ форматирует библиографию по ГОСТ, исправляет стиль и создаёт аннотацию</div>', unsafe_allow_html=True)

        input_method = st.radio("Способ ввода", ["✏️ Вставить текст", "📁 Загрузить файл (.txt)"], horizontal=True)

        raw_text = ""
        if "Вставить" in input_method:
            raw_text = st.text_area("Текст для обработки", height=250,
                                    placeholder="Вставьте текст реферата, список литературы или черновик...")
        else:
            uploaded = st.file_uploader("Загрузите .txt файл", type=["txt"])
            if uploaded:
                raw_text = uploaded.read().decode("utf-8", errors="replace")
                st.info(f"📄 Загружен файл: {uploaded.name} ({len(raw_text)} символов)")

        col_a, col_b = st.columns(2)
        with col_a:
            do_gost     = st.checkbox("📚 Форматировать список литературы по ГОСТ", value=True)
            do_style    = st.checkbox("✍️ Исправить стилистические ошибки", value=True)
        with col_b:
            do_abstract = st.checkbox("📄 Создать аннотацию (абстракт)", value=True)
            do_plan     = st.checkbox("🗂 Составить план работы", value=True)

        correct_btn = st.button("🔍 Обработать текст", use_container_width=True)

        if correct_btn:
            if not raw_text.strip():
                st.warning("⚠️ Введите или загрузите текст.")
            else:
                tasks = []
                if do_gost:     tasks.append("1. Найди все источники в тексте и отформатируй их по ГОСТ 7.0.5-2008. Если список литературы уже есть — исправь его. Выведи раздел «СПИСОК ЛИТЕРАТУРЫ (ГОСТ)».")
                if do_style:    tasks.append("2. Исправь грубые стилистические, лексические и грамматические ошибки. Выведи исправленную версию текста в разделе «ИСПРАВЛЕННЫЙ ТЕКСТ».")
                if do_abstract: tasks.append("3. Напиши краткую аннотацию (абстракт, 80–120 слов) в разделе «АННОТАЦИЯ».")
                if do_plan:     tasks.append("4. Составь структурированный план работы (введение, главы, заключение) в разделе «ПЛАН РАБОТЫ».")

                if not tasks:
                    st.warning("⚠️ Выберите хотя бы одну задачу.")
                else:
                    prompt = f"""Ты — академический редактор. Обработай следующий текст и выполни задачи ниже.

ТЕКСТ:
\"\"\"{raw_text[:6000]}\"\"\"

ЗАДАЧИ:
{chr(10).join(tasks)}

Формат: чёткие разделы с заголовками, как указано в каждой задаче. Не добавляй ничего лишнего."""

                    with st.spinner("🤖 Обрабатываю текст..."):
                        result = call_gemini(prompt)

                    if result.startswith("__ERROR__"):
                        st.error(result.replace("__ERROR__: ", ""))
                    else:
                        st.markdown("### 📋 Результат обработки")
                        st.markdown(f'<div class="result-box">{result}</div>', unsafe_allow_html=True)
                        st.download_button(
                            "⬇️ Скачать результат (.txt)",
                            data=result.encode("utf-8"),
                            file_name="gost_result.txt",
                            mime="text/plain",
                        )

# ─────────────────────────────────────────
# ── ПРЕПОДАВАТЕЛЬ ─────────────────────────
# ─────────────────────────────────────────
else:
    tab3, tab4 = st.tabs(["❓ Квиз-Конструктор", "📊 Ассистент Оценки"])

    # ── TAB 3: КВИЗ ────────────────────────
    with tab3:
        st.markdown('<div class="module-card"><b>❓ Квиз-Конструктор</b><br>Загрузите материал или введите тему — ИИ составит тест с ответами и объяснениями</div>', unsafe_allow_html=True)

        q_input_method = st.radio("Источник материала", ["✏️ Ввести тему", "📁 Загрузить текст (.txt)"], horizontal=True)

        q_source = ""
        if "Ввести тему" in q_input_method:
            q_source = st.text_input("Тема теста", placeholder="Например: Клеточное дыхание, Фотосинтез, Второй закон Ньютона...")
        else:
            q_file = st.file_uploader("Загрузите лекцию / главу учебника (.txt)", type=["txt"])
            if q_file:
                q_source = q_file.read().decode("utf-8", errors="replace")
                st.info(f"📄 {q_file.name} ({len(q_source)} символов)")

        col1, col2 = st.columns(2)
        with col1:
            q_count = st.selectbox("Количество вопросов", [5, 10, 15])
        with col2:
            q_type = st.selectbox("Тип вопросов", ["Один правильный ответ (A/B/C/D)", "Открытые вопросы"])

        q_lang = st.selectbox("Язык теста", ["Русский", "English", "Uzbek"])
        quiz_btn = st.button("📝 Сгенерировать тест", use_container_width=True)

        if quiz_btn:
            if not q_source.strip():
                st.warning("⚠️ Введите тему или загрузите материал.")
            else:
                if "Один правильный" in q_type:
                    fmt_instruction = """For each question provide:
- Question text
- 4 answer options labeled A), B), C), D)
- Correct answer label
- Explanation (2-3 sentences why this answer is correct)

Format each question exactly like this:
**Вопрос N:** [question text]
A) [option]
B) [option]
C) [option]
D) [option]
✅ Правильный ответ: [letter]) [text]
💡 Объяснение: [explanation]
---"""
                else:
                    fmt_instruction = """For each question provide:
- Open-ended question
- Model answer (3-5 sentences)
- Key evaluation criteria

Format each question exactly like this:
**Вопрос N:** [question text]
📝 Эталонный ответ: [model answer]
🔑 Критерии оценки: [criteria]
---"""

                prompt = f"""You are a professional educator. Create a quiz with exactly {q_count} questions on the following material.
Language: {q_lang}
Source: \"\"\"{q_source[:5000]}\"\"\"

{fmt_instruction}

Generate exactly {q_count} questions. Cover different aspects of the material."""

                with st.spinner("🤖 Генерирую тест..."):
                    result = call_gemini(prompt)

                if result.startswith("__ERROR__"):
                    st.error(result.replace("__ERROR__: ", ""))
                else:
                    st.markdown("### 📋 Готовый тест")
                    st.markdown(f'<div class="result-box">{result}</div>', unsafe_allow_html=True)

                    col_dl1, col_dl2 = st.columns(2)
                    with col_dl1:
                        st.download_button(
                            "⬇️ Скачать TXT",
                            data=result.encode("utf-8"),
                            file_name="quiz.txt",
                            mime="text/plain",
                            use_container_width=True,
                        )
                    with col_dl2:
                        md_content = f"# Тест\n\n{result}"
                        st.download_button(
                            "⬇️ Скачать Markdown",
                            data=md_content.encode("utf-8"),
                            file_name="quiz.md",
                            mime="text/markdown",
                            use_container_width=True,
                        )

    # ── TAB 4: ОЦЕНКА ──────────────────────
    with tab4:
        st.markdown('<div class="module-card"><b>📊 Ассистент Оценки и Фидбека</b><br>ИИ анализирует студенческую работу и пишет развёрнутый отзыв с рекомендательной оценкой</div>', unsafe_allow_html=True)

        essay_text = st.text_area(
            "Текст студенческой работы / эссе / ответа",
            height=220,
            placeholder="Вставьте текст работы студента...",
        )

        criteria = st.text_area(
            "Критерии оценки (необязательно)",
            height=100,
            placeholder="Например: раскрытие темы, использование терминов, наличие примеров, логика изложения, грамотность...",
        )

        col1, col2 = st.columns(2)
        with col1:
            subject = st.text_input("Предмет / дисциплина", placeholder="Биология, История, Математика...")
            max_grade = st.selectbox("Максимальный балл", [5, 10, 20, 100])
        with col2:
            grade_system = st.selectbox("Система оценивания", ["Пятибалльная (1–5)", "По 10-балльной шкале", "По 100-балльной шкале", "Pass/Fail"])
            feedback_lang = st.selectbox("Язык фидбека", ["Русский", "English", "Uzbek"])

        eval_btn = st.button("🔍 Оценить работу", use_container_width=True)

        if eval_btn:
            if not essay_text.strip():
                st.warning("⚠️ Вставьте текст студенческой работы.")
            else:
                criteria_block = f"Критерии оценки:\n{criteria}" if criteria.strip() else "Оцени по общим академическим критериям."

                prompt = f"""Ты — опытный преподаватель по дисциплине «{subject or 'общеакадемическая дисциплина'}».
Проанализируй студенческую работу ниже и составь развёрнутый педагогический отзыв.
Язык отзыва: {feedback_lang}.
Система оценивания: {grade_system} (максимум: {max_grade}).
{criteria_block}

РАБОТА СТУДЕНТА:
\"\"\"{essay_text[:6000]}\"\"\"

Структура твоего ответа:

## 🌟 Сильные стороны
[Перечисли 3–5 конкретных достоинств работы]

## ⚠️ Области для улучшения
[Перечисли 3–5 конкретных недостатков с пояснениями]

## 📌 Конкретные рекомендации
[Дай 3–4 практических совета, как улучшить работу]

## 🎯 Рекомендательная оценка
[Выставь оценку: X из {max_grade}. Обоснуй в 2–3 предложениях]

## 💬 Итоговый отзыв для студента
[Напиши вежливый, мотивирующий отзыв в 4–6 предложениях, адресованный непосредственно студенту]"""

                with st.spinner("🤖 Анализирую работу..."):
                    result = call_gemini(prompt)

                if result.startswith("__ERROR__"):
                    st.error(result.replace("__ERROR__: ", ""))
                else:
                    st.markdown("### 📋 Результаты оценки")
                    st.markdown(result)
                    st.divider()
                    st.download_button(
                        "⬇️ Скачать отзыв (.txt)",
                        data=result.encode("utf-8"),
                        file_name="feedback.txt",
                        mime="text/plain",
                    )

# ─────────────────────────────────────────
# FOOTER
# ─────────────────────────────────────────
st.divider()
st.caption("🎓 EduHelper AI · Бесплатный академический помощник · Google Gemini 1.5 Flash · Создано с ❤️ на Streamlit")
